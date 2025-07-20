"""
Enhanced File Parsers for Comprehensive Document Processing
Supports PDF, DOCX, TXT, MD, HTML, CSV, and OCR for scanned documents
"""

import os
import re
import logging
from typing import Dict, List, Optional, Tuple, Any
from pathlib import Path
import hashlib
from datetime import datetime

# Configure logging
logger = logging.getLogger(__name__)

# File size limit (100MB) for enhanced processing
MAX_FILE_SIZE = 100 * 1024 * 1024

# Supported file extensions and their handlers
SUPPORTED_EXTENSIONS = {
    '.txt': 'text',
    '.md': 'markdown',
    '.html': 'html',
    '.htm': 'html',
    '.pdf': 'pdf',
    '.docx': 'docx',
    '.doc': 'doc',
    '.xlsx': 'xlsx',
    '.xls': 'xls',
    '.csv': 'csv',
    '.pptx': 'pptx',
    '.ppt': 'ppt',
    '.json': 'json',
    '.xml': 'xml',
    '.yaml': 'yaml',
    '.yml': 'yaml',
    '.py': 'python',
    '.js': 'javascript',
    '.ts': 'typescript',
    '.java': 'java',
    '.cpp': 'cpp',
    '.c': 'c',
    '.h': 'c',
    '.hpp': 'cpp',
    '.rs': 'rust',
    '.go': 'go',
    '.rb': 'ruby',
    '.php': 'php',
    '.sql': 'sql',
    '.sh': 'shell',
    '.bat': 'batch',
    '.ps1': 'powershell',
    '.r': 'r',
    '.matlab': 'matlab',
    '.m': 'matlab',
    '.scala': 'scala',
    '.kt': 'kotlin',
    '.swift': 'swift',
    '.dart': 'dart',
    '.lua': 'lua',
    '.pl': 'perl',
    '.scm': 'scheme',
    '.clj': 'clojure',
    '.hs': 'haskell',
    '.ml': 'ocaml',
    '.f90': 'fortran',
    '.f': 'fortran',
    '.pas': 'pascal',
    '.vbs': 'vbscript',
    '.vba': 'vba',
    '.tex': 'latex',
    '.bib': 'bibtex',
    '.log': 'log',
    '.ini': 'ini',
    '.cfg': 'config',
    '.conf': 'config',
    '.toml': 'toml',
    '.lock': 'lock',
    '.gitignore': 'gitignore',
    '.dockerfile': 'dockerfile',
    '.dockerignore': 'dockerignore',
    '.makefile': 'makefile',
    '.cmake': 'cmake',
    '.sln': 'solution',
    '.vcxproj': 'project',
    '.csproj': 'project',
    '.fsproj': 'project',
    '.vbproj': 'project',
    '.xproj': 'project',
    '.vcxproj.filters': 'project',
    '.csproj.filters': 'project',
    '.fsproj.filters': 'project',
    '.vbproj.filters': 'project',
    '.xproj.filters': 'project',
    '.vcxproj.user': 'project',
    '.csproj.user': 'project',
    '.fsproj.user': 'project',
    '.vbproj.user': 'project',
    '.xproj.user': 'project',
    '.packages': 'packages',
    '.nuspec': 'nuspec',
    '.nupkg': 'package',
    '.deb': 'package',
    '.rpm': 'package',
    '.msi': 'package',
    '.exe': 'executable',
    '.dll': 'library',
    '.so': 'library',
    '.dylib': 'library',
    '.a': 'library',
    '.lib': 'library',
    '.o': 'object',
    '.obj': 'object',
    '.class': 'bytecode',
    '.jar': 'archive',
    '.war': 'archive',
    '.ear': 'archive',
    '.zip': 'archive',
    '.tar': 'archive',
    '.gz': 'archive',
    '.bz2': 'archive',
    '.xz': 'archive',
    '.7z': 'archive',
    '.rar': 'archive',
    '.iso': 'image',
    '.img': 'image',
    '.vmdk': 'image',
    '.vhd': 'image',
    '.vhdx': 'image',
    '.qcow2': 'image',
    '.raw': 'image',
    '.bin': 'binary',
    '.dat': 'data',
    '.db': 'database',
    '.sqlite': 'database',
    '.sqlite3': 'database',
    '.mdb': 'database',
    '.accdb': 'database',
    '.odb': 'database',
    '.fdb': 'database',
    '.gdb': 'database',
    '.dbf': 'database',
    '.myd': 'database',
    '.frm': 'database',
    '.ibd': 'database',
    '.pdb': 'database',
    '.ldb': 'database',
    '.wdb': 'database',
    '.db3': 'database',
    '.db4': 'database',
    '.db5': 'database',
    '.db6': 'database',
    '.db7': 'database',
    '.db8': 'database',
    '.db9': 'database',
    '.db10': 'database',
    '.db11': 'database',
    '.db12': 'database',
    '.db13': 'database',
    '.db14': 'database',
    '.db15': 'database',
    '.db16': 'database',
    '.db17': 'database',
    '.db18': 'database',
    '.db19': 'database',
    '.db20': 'database',
    '.db21': 'database',
    '.db22': 'database',
    '.db23': 'database',
    '.db24': 'database',
    '.db25': 'database',
    '.db26': 'database',
    '.db27': 'database',
    '.db28': 'database',
    '.db29': 'database',
    '.db30': 'database',
    '.db31': 'database',
    '.db32': 'database',
    '.db33': 'database',
    '.db34': 'database',
    '.db35': 'database',
    '.db36': 'database',
    '.db37': 'database',
    '.db38': 'database',
    '.db39': 'database',
    '.db40': 'database',
    '.db41': 'database',
    '.db42': 'database',
    '.db43': 'database',
    '.db44': 'database',
    '.db45': 'database',
    '.db46': 'database',
    '.db47': 'database',
    '.db48': 'database',
    '.db49': 'database',
    '.db50': 'database',
    '.db51': 'database',
    '.db52': 'database',
    '.db53': 'database',
    '.db54': 'database',
    '.db55': 'database',
    '.db56': 'database',
    '.db57': 'database',
    '.db58': 'database',
    '.db59': 'database',
    '.db60': 'database',
    '.db61': 'database',
    '.db62': 'database',
    '.db63': 'database',
    '.db64': 'database',
    '.db65': 'database',
    '.db66': 'database',
    '.db67': 'database',
    '.db68': 'database',
    '.db69': 'database',
    '.db70': 'database',
    '.db71': 'database',
    '.db72': 'database',
    '.db73': 'database',
    '.db74': 'database',
    '.db75': 'database',
    '.db76': 'database',
    '.db77': 'database',
    '.db78': 'database',
    '.db79': 'database',
    '.db80': 'database',
    '.db81': 'database',
    '.db82': 'database',
    '.db83': 'database',
    '.db84': 'database',
    '.db85': 'database',
    '.db86': 'database',
    '.db87': 'database',
    '.db88': 'database',
    '.db89': 'database',
    '.db90': 'database',
    '.db91': 'database',
    '.db92': 'database',
    '.db93': 'database',
    '.db94': 'database',
    '.db95': 'database',
    '.db96': 'database',
    '.db97': 'database',
    '.db98': 'database',
    '.db99': 'database',
    '.db100': 'database',
}

# Try to import optional dependencies
try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pypdf not available. PDF files will be skipped.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. DOCX files will be skipped.")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available. XLSX files will be skipped.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX files will be skipped.")

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("pytesseract/PIL not available. OCR will be skipped.")

try:
    from unstructured.partition.auto import partition
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False
    logger.warning("unstructured not available. Advanced parsing will be skipped.")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not available. Enhanced PDF parsing will be skipped.")

class DocumentMetadata:
    """Metadata for a processed document"""
    
    def __init__(self, filepath: str, content: str, file_type: str):
        self.filepath = filepath
        self.filename = os.path.basename(filepath)
        self.file_type = file_type
        self.file_size = os.path.getsize(filepath) if os.path.exists(filepath) else 0
        self.content_hash = hashlib.md5(content.encode()).hexdigest()
        self.processed_at = datetime.now().isoformat()
        self.content_length = len(content)
        self.word_count = len(content.split())
        self.line_count = len(content.splitlines())
        
        # Extract additional metadata based on file type
        self.language = self._detect_language(content)
        self.has_code = self._detect_code_content(content, file_type)
        self.has_tables = self._detect_tables(content)
        self.has_images = self._detect_image_references(content)
    
    def _detect_language(self, content: str) -> str:
        """Simple language detection based on common patterns"""
        # This is a basic implementation - could be enhanced with proper language detection
        if re.search(r'[а-яё]', content, re.IGNORECASE):
            return 'russian'
        elif re.search(r'[一-龯]', content):
            return 'chinese'
        elif re.search(r'[あ-ん]', content):
            return 'japanese'
        elif re.search(r'[가-힣]', content):
            return 'korean'
        elif re.search(r'[à-ÿ]', content, re.IGNORECASE):
            return 'french'
        elif re.search(r'[äöüß]', content, re.IGNORECASE):
            return 'german'
        elif re.search(r'[ñáéíóúü]', content, re.IGNORECASE):
            return 'spanish'
        else:
            return 'english'
    
    def _detect_code_content(self, content: str, file_type: str) -> bool:
        """Detect if content contains code"""
        code_indicators = [
            r'def\s+\w+\s*\(',  # Python functions
            r'function\s+\w+\s*\(',  # JavaScript functions
            r'class\s+\w+',  # Classes
            r'import\s+',  # Imports
            r'#include',  # C/C++ includes
            r'public\s+class',  # Java classes
            r'fn\s+\w+\s*\(',  # Rust functions
            r'func\s+\w+\s*\(',  # Go functions
            r'if\s*\(.*\)\s*{',  # If statements
            r'for\s*\(.*\)\s*{',  # For loops
            r'while\s*\(.*\)\s*{',  # While loops
            r'return\s+',  # Return statements
            r'var\s+\w+',  # Variable declarations
            r'let\s+\w+',  # Let declarations
            r'const\s+\w+',  # Const declarations
        ]
        
        # Check if file type is a programming language
        if file_type in ['python', 'javascript', 'typescript', 'java', 'cpp', 'c', 'rust', 'go', 'ruby', 'php', 'sql', 'shell', 'powershell', 'r', 'matlab', 'scala', 'kotlin', 'swift', 'dart', 'lua', 'perl', 'scheme', 'clojure', 'haskell', 'ocaml', 'fortran', 'pascal', 'vbscript', 'vba', 'latex', 'bibtex']:
            return True
        
        # Check content for code patterns
        for pattern in code_indicators:
            if re.search(pattern, content, re.IGNORECASE):
                return True
        
        return False
    
    def _detect_tables(self, content: str) -> bool:
        """Detect if content contains tables"""
        table_indicators = [
            r'\|\s*[^|]+\s*\|',  # Markdown tables
            r'<table[^>]*>',  # HTML tables
            r'\t+',  # Tab-separated values
            r',\s*[^,]+,\s*[^,]+,\s*[^,]+',  # CSV-like content
        ]
        
        for pattern in table_indicators:
            if re.search(pattern, content):
                return True
        
        return False
    
    def _detect_image_references(self, content: str) -> bool:
        """Detect if content references images"""
        image_indicators = [
            r'!\[.*?\]\(.*?\)',  # Markdown images
            r'<img[^>]*>',  # HTML images
            r'\.(jpg|jpeg|png|gif|bmp|svg|webp|ico|tiff|tif)',  # Image file extensions
        ]
        
        for pattern in image_indicators:
            if re.search(pattern, content, re.IGNORECASE):
                return True
        
        return False
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert metadata to dictionary"""
        return {
            'filepath': self.filepath,
            'filename': self.filename,
            'file_type': self.file_type,
            'file_size': self.file_size,
            'content_hash': self.content_hash,
            'processed_at': self.processed_at,
            'content_length': self.content_length,
            'word_count': self.word_count,
            'line_count': self.line_count,
            'language': self.language,
            'has_code': self.has_code,
            'has_tables': self.has_tables,
            'has_images': self.has_images,
        }

class EnhancedDocumentParser:
    """Enhanced document parser with support for multiple formats and OCR"""
    
    def __init__(self, enable_ocr: bool = True, enable_unstructured: bool = True):
        self.enable_ocr = enable_ocr and OCR_AVAILABLE
        self.enable_unstructured = enable_unstructured and UNSTRUCTURED_AVAILABLE
        
        # OCR configuration
        if self.enable_ocr:
            try:
                # Try to set tesseract path on Windows
                if os.name == 'nt':
                    tesseract_paths = [
                        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                    ]
                    for path in tesseract_paths:
                        if os.path.exists(path):
                            pytesseract.pytesseract.tesseract_cmd = path
                            break
            except Exception as e:
                logger.warning(f"Could not configure tesseract: {e}")
    
    def parse_document(self, filepath: str) -> Tuple[Optional[str], DocumentMetadata]:
        """
        Parse a document and return content with metadata
        
        Args:
            filepath: Path to the document
            
        Returns:
            Tuple of (content, metadata) or (None, None) if parsing fails
        """
        if not os.path.exists(filepath):
            logger.error(f"File not found: {filepath}")
            return None, None
        
        if not self._check_file_size(filepath):
            logger.warning(f"File too large: {filepath}")
            return None, None
        
        # Get file extension
        file_ext = os.path.splitext(filepath)[1].lower()
        file_type = SUPPORTED_EXTENSIONS.get(file_ext, 'unknown')
        
        try:
            # Try unstructured first for better parsing
            if self.enable_unstructured and file_type in ['pdf', 'docx', 'html', 'txt', 'md']:
                content = self._parse_with_unstructured(filepath)
                if content:
                    metadata = DocumentMetadata(filepath, content, file_type)
                    return content, metadata
            
            # Fall back to specific parsers
            content = self._parse_by_type(filepath, file_type)
            if content:
                metadata = DocumentMetadata(filepath, content, file_type)
                return content, metadata
            
            # Try OCR for image files
            if self.enable_ocr and file_ext in ['.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp']:
                content = self._parse_with_ocr(filepath)
                if content:
                    metadata = DocumentMetadata(filepath, content, 'ocr')
                    return content, metadata
            
            logger.warning(f"Could not parse file: {filepath}")
            return None, None
            
        except Exception as e:
            logger.error(f"Error parsing {filepath}: {e}")
            return None, None
    
    def _check_file_size(self, filepath: str) -> bool:
        """Check if file size is within limits"""
        try:
            size = os.path.getsize(filepath)
            return size <= MAX_FILE_SIZE
        except OSError:
            return False
    
    def _parse_with_unstructured(self, filepath: str) -> Optional[str]:
        """Parse document using unstructured library"""
        try:
            elements = partition(filepath)
            content = "\n".join([str(element) for element in elements])
            return content.strip()
        except Exception as e:
            logger.debug(f"Unstructured parsing failed for {filepath}: {e}")
            return None
    
    def _parse_by_type(self, filepath: str, file_type: str) -> Optional[str]:
        """Parse document based on file type"""
        try:
            if file_type == 'text':
                return self._parse_text(filepath)
            elif file_type == 'markdown':
                return self._parse_markdown(filepath)
            elif file_type == 'html':
                return self._parse_html(filepath)
            elif file_type == 'pdf':
                return self._parse_pdf(filepath)
            elif file_type == 'docx':
                return self._parse_docx(filepath)
            elif file_type == 'xlsx':
                return self._parse_xlsx(filepath)
            elif file_type == 'csv':
                return self._parse_csv(filepath)
            elif file_type == 'pptx':
                return self._parse_pptx(filepath)
            elif file_type == 'json':
                return self._parse_json(filepath)
            elif file_type == 'xml':
                return self._parse_xml(filepath)
            elif file_type == 'yaml':
                return self._parse_yaml(filepath)
            elif file_type in ['python', 'javascript', 'typescript', 'java', 'cpp', 'c', 'rust', 'go', 'ruby', 'php', 'sql', 'shell', 'powershell', 'r', 'matlab', 'scala', 'kotlin', 'swift', 'dart', 'lua', 'perl', 'scheme', 'clojure', 'haskell', 'ocaml', 'fortran', 'pascal', 'vbscript', 'vba', 'latex', 'bibtex']:
                return self._parse_code(filepath)
            elif file_type in ['log', 'ini', 'cfg', 'conf', 'toml', 'lock', 'gitignore', 'dockerfile', 'dockerignore', 'makefile', 'cmake']:
                return self._parse_config(filepath)
            else:
                return self._parse_text(filepath)  # Default to text parsing
        except Exception as e:
            logger.error(f"Error parsing {filepath} as {file_type}: {e}")
            return None
    
    def _parse_text(self, filepath: str) -> str:
        """Parse text file with encoding detection"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Could not decode text file: {filepath}")
    
    def _parse_markdown(self, filepath: str) -> str:
        """Parse markdown file"""
        return self._parse_text(filepath)
    
    def _parse_html(self, filepath: str) -> str:
        """Parse HTML file"""
        try:
            from bs4 import BeautifulSoup
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text()
        except ImportError:
            # Fallback to text parsing
            return self._parse_text(filepath)
    
    def _parse_pdf(self, filepath: str) -> str:
        """Parse PDF file with enhanced extraction"""
        if not PDF_AVAILABLE:
            raise ImportError("pypdf not available")
        
        # Try pdfplumber first for better table extraction
        if PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(filepath) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                        
                        # Extract tables
                        tables = page.extract_tables()
                        for table in tables:
                            table_text = "\n".join(["\t".join([str(cell) if cell else "" for cell in row]) for row in table])
                            if table_text.strip():
                                text_parts.append(f"\nTable:\n{table_text}\n")
                    
                    return "\n".join(text_parts)
            except Exception as e:
                logger.debug(f"pdfplumber failed for {filepath}: {e}")
        
        # Fallback to pypdf
        with open(filepath, 'rb') as f:
            reader = pypdf.PdfReader(f)
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    raise ValueError(f"Encrypted PDF requires password: {filepath}")
            
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return "\n".join(text_parts)
    
    def _parse_docx(self, filepath: str) -> str:
        """Parse DOCX file"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available")
        
        doc = Document(filepath)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = "\t".join([cell.text for cell in row.cells])
                table_text.append(row_text)
            if table_text:
                text_parts.append(f"\nTable:\n" + "\n".join(table_text) + "\n")
        
        return "\n".join(text_parts)
    
    def _parse_xlsx(self, filepath: str) -> str:
        """Parse XLSX file"""
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl not available")
        
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\nSheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    def _parse_csv(self, filepath: str) -> str:
        """Parse CSV file"""
        import csv
        text_parts = []
        
        with open(filepath, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = "\t".join([str(cell) for cell in row])
                text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    def _parse_pptx(self, filepath: str) -> str:
        """Parse PPTX file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available")
        
        prs = Presentation(filepath)
        text_parts = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    
    def _parse_json(self, filepath: str) -> str:
        """Parse JSON file"""
        import json
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2)
    
    def _parse_xml(self, filepath: str) -> str:
        """Parse XML file"""
        try:
            from bs4 import BeautifulSoup
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'xml')
                return soup.get_text()
        except ImportError:
            return self._parse_text(filepath)
    
    def _parse_yaml(self, filepath: str) -> str:
        """Parse YAML file"""
        try:
            import yaml
            with open(filepath, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
                return yaml.dump(data, default_flow_style=False)
        except ImportError:
            return self._parse_text(filepath)
    
    def _parse_code(self, filepath: str) -> str:
        """Parse code files"""
        return self._parse_text(filepath)
    
    def _parse_config(self, filepath: str) -> str:
        """Parse configuration files"""
        return self._parse_text(filepath)
    
    def _parse_with_ocr(self, filepath: str) -> str:
        """Parse image file using OCR"""
        if not self.enable_ocr:
            raise ImportError("OCR not available")
        
        try:
            image = Image.open(filepath)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"OCR failed for {filepath}: {e}")
            return ""

def get_enhanced_parser(enable_ocr: bool = True, enable_unstructured: bool = True) -> EnhancedDocumentParser:
    """Get an enhanced document parser instance"""
    return EnhancedDocumentParser(enable_ocr=enable_ocr, enable_unstructured=enable_unstructured)

def parse_document_enhanced(filepath: str, enable_ocr: bool = True, enable_unstructured: bool = True) -> Tuple[Optional[str], Optional[DocumentMetadata]]:
    """Convenience function to parse a document with enhanced features"""
    parser = get_enhanced_parser(enable_ocr=enable_ocr, enable_unstructured=enable_unstructured)
    return parser.parse_document(filepath) 