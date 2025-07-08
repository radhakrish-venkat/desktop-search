import os
import sys
import re
from typing import Tuple, Optional, Union
import logging

# Configure logging
logging.basicConfig(level=logging.WARNING)
logger = logging.getLogger(__name__)

# File size limit (50MB) to prevent memory issues
MAX_FILE_SIZE = 50 * 1024 * 1024

# Suspicious patterns to detect in file content
SUSPICIOUS_PATTERNS = [
    r'<script[^>]*>.*?</script>',  # Script tags
    r'javascript:',                # JavaScript protocol
    r'data:text/html',            # Data URLs
    r'vbscript:',                 # VBScript
    r'<iframe[^>]*>',             # Iframe tags
    r'onload\s*=',                # Event handlers
    r'onerror\s*=',               # Event handlers
    r'onclick\s*=',               # Event handlers
    r'<object[^>]*>',             # Object tags
    r'<embed[^>]*>',              # Embed tags
    r'<applet[^>]*>',             # Applet tags
    r'<meta[^>]*http-equiv[^>]*refresh[^>]*>',  # Meta refresh
    r'<link[^>]*javascript[^>]*>',  # JavaScript links
    r'<a[^>]*href\s*=\s*["\']javascript:',  # JavaScript links
]

# Supported MIME types for file processing
SUPPORTED_MIME_TYPES = {
    'text/plain': '.txt',
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'application/msword': '.doc',
    'application/vnd.ms-excel': '.xls',
    'application/vnd.ms-powerpoint': '.ppt',
    'text/csv': '.csv',
    'application/csv': '.csv',
}

try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pypdf not available. PDF files will be skipped.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. DOCX files will be skipped.")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available. XLSX files will be skipped.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX files will be skipped.")

def _check_file_size(filepath: str) -> bool:
    """Check if file size is within acceptable limits."""
    try:
        size = os.path.getsize(filepath)
        return size <= MAX_FILE_SIZE
    except OSError:
        return False

def _validate_file_content(filepath: str, content: str) -> bool:
    """
    Validate file content for suspicious patterns.
    
    Args:
        filepath: Path to the file being validated
        content: File content to validate
        
    Returns:
        True if content is safe, False if suspicious patterns detected
    """
    if not content:
        return True
    
    for pattern in SUSPICIOUS_PATTERNS:
        if re.search(pattern, content, re.IGNORECASE | re.DOTALL):
            logger.warning(f"Suspicious content detected in {filepath}: {pattern}")
            return False
    
    return True

def _get_file_mime_type(filepath: str) -> Optional[str]:
    """
    Get MIME type of file using python-magic if available.
    
    Args:
        filepath: Path to the file
        
    Returns:
        MIME type string or None if detection fails
    """
    try:
        import magic
        return magic.from_file(filepath, mime=True)
    except ImportError:
        # Fallback to file extension if python-magic not available
        logger.debug("python-magic not available, using file extension for MIME type detection")
        return None
    except Exception as e:
        logger.debug(f"Error detecting MIME type for {filepath}: {e}")
        return None

def _is_supported_mime_type(mime_type: str) -> bool:
    """
    Check if MIME type is supported for processing.
    
    Args:
        mime_type: MIME type to check
        
    Returns:
        True if supported, False otherwise
    """
    return mime_type in SUPPORTED_MIME_TYPES

def get_text_from_txt(filepath: str) -> str:
    """
    Extracts text from a plain text file.
    
    Args:
        filepath: Path to the text file
        
    Returns:
        Extracted text as string
    """
    if not _check_file_size(filepath):
        logger.warning(f"File too large: {filepath}")
        return ""
        
    try:
        # Try UTF-8 first, then fallback to other encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    content = f.read()
                    
                    # Validate content for suspicious patterns
                    if not _validate_file_content(filepath, content):
                        logger.warning(f"Content validation failed for {filepath}")
                        return ""
                    
                    return content
            except UnicodeDecodeError:
                continue
        logger.error(f"Could not decode text file: {filepath}")
        return ""
    except Exception as e:
        logger.error(f"Error reading TXT file {filepath}: {e}")
        return ""

def get_text_from_pdf(filepath: str) -> str:
    """
    Extracts text from a PDF file.
    
    Args:
        filepath: Path to the PDF file
        
    Returns:
        Extracted text as string
    """
    if not PDF_AVAILABLE:
        logger.warning(f"PDF support not available, skipping: {filepath}")
        return ""
        
    if not _check_file_size(filepath):
        logger.warning(f"PDF file too large: {filepath}")
        return ""
        
    text = ""
    try:
        with open(filepath, 'rb') as f:
            reader = pypdf.PdfReader(f)
            
            # Handle encrypted PDFs
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    logger.warning(f"Encrypted PDF requires password: {filepath}")
                    return ""

            # Extract text from all pages
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num} in {filepath}: {e}")
                    continue
                    
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading PDF file {filepath}: {e}")
        return ""

def get_text_from_docx(filepath: str) -> str:
    """
    Extracts text from a Word document (.docx).
    
    Args:
        filepath: Path to the DOCX file
        
    Returns:
        Extracted text as string
    """
    if not DOCX_AVAILABLE:
        logger.warning(f"DOCX support not available, skipping: {filepath}")
        return ""
        
    if not _check_file_size(filepath):
        logger.warning(f"DOCX file too large: {filepath}")
        return ""
        
    text = ""
    try:
        document = Document(filepath)
        
        # Extract text from paragraphs
        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Remove trailing newline if present
        if text.endswith('\n'):
            text = text[:-1]
        
        # Extract text from tables
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            text += paragraph.text + "\n"
                            
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading DOCX file {filepath}: {e}")
        return ""

def get_text_from_xlsx(filepath: str) -> str:
    """
    Extracts text from an Excel workbook (.xlsx).
    
    Args:
        filepath: Path to the XLSX file
        
    Returns:
        Extracted text as string
    """
    if not XLSX_AVAILABLE:
        logger.warning(f"XLSX support not available, skipping: {filepath}")
        return ""
        
    if not _check_file_size(filepath):
        logger.warning(f"XLSX file too large: {filepath}")
        return ""
        
    text = ""
    try:
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"--- Sheet: {sheet_name} ---\n"
            
            # Get the used range
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value is not None:
                        row_text.append(str(cell.value))
                if row_text:
                    text += " ".join(row_text) + "\n"
                    
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading XLSX file {filepath}: {e}")
        return ""

def get_text_from_pptx(filepath: str) -> str:
    """
    Extracts text from a PowerPoint presentation (.pptx).
    
    Args:
        filepath: Path to the PPTX file
        
    Returns:
        Extracted text as string
    """
    if not PPTX_AVAILABLE:
        logger.warning(f"PPTX support not available, skipping: {filepath}")
        return ""
        
    if not _check_file_size(filepath):
        logger.warning(f"PPTX file too large: {filepath}")
        return ""
        
    text = ""
    try:
        prs = Presentation(filepath)
        
        for slide_num, slide in enumerate(prs.slides):
            text += f"--- Slide {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                # Extract text from text frames
                if hasattr(shape, "text_frame") and shape.text_frame:  # type: ignore
                    text_frame = getattr(shape, "text_frame")
                    if text_frame.text.strip():
                        text += text_frame.text + "\n"
                
                # Extract text from tables
                elif hasattr(shape, "table") and shape.table:  # type: ignore
                    table = getattr(shape, "table")
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    text += paragraph.text + "\n"
                                    
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading PPTX file {filepath}: {e}")
        return ""

def get_text_from_csv(filepath: str) -> str:
    """
    Extracts text from a CSV file.
    
    Args:
        filepath: Path to the CSV file
        
    Returns:
        Extracted text as string
    """
    if not _check_file_size(filepath):
        logger.warning(f"CSV file too large: {filepath}")
        return ""
        
    text = ""
    try:
        import csv
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    for row_num, row in enumerate(reader, 1):
                        if row:  # Skip empty rows
                            text += f"Row {row_num}: {' | '.join(str(cell) for cell in row)}\n"
                break  # If successful, break out of encoding loop
            except UnicodeDecodeError:
                continue
        else:
            logger.error(f"Could not decode CSV file: {filepath}")
            return ""
            
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading CSV file {filepath}: {e}")
        return ""

def get_text_from_doc(filepath: str) -> str:
    """
    Extracts text from a legacy Word document (.doc).
    
    Args:
        filepath: Path to the DOC file
        
    Returns:
        Extracted text as string
    """
    if not _check_file_size(filepath):
        logger.warning(f"DOC file too large: {filepath}")
        return ""
        
    text = ""
    try:
        # Try using python-docx2txt for .doc files
        try:
            import docx2txt
            text = docx2txt.process(filepath)
            return text.strip()
        except ImportError:
            logger.warning("docx2txt not available for .doc files, trying alternative method")
        
        # Alternative: Try using antiword if available
        try:
            import subprocess
            result = subprocess.run(['antiword', filepath], capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                return result.stdout.strip()
        except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
            logger.warning("antiword not available for .doc files")
        
        # Fallback: Try to extract as binary and look for text patterns
        try:
            with open(filepath, 'rb') as f:
                content = f.read()
                # Look for text patterns in binary content
                text_patterns = re.findall(rb'[\x20-\x7E]{4,}', content)
                text = ' '.join(pattern.decode('utf-8', errors='ignore') for pattern in text_patterns)
                return text.strip()
        except Exception:
            pass
            
        logger.warning(f"Could not extract text from DOC file: {filepath}")
        return ""
        
    except Exception as e:
        logger.error(f"Error reading DOC file {filepath}: {e}")
        return ""

def get_text_from_file(filepath: str) -> Tuple[Optional[str], str]:
    """
    Dispatches to the correct text extraction function based on file extension.
    Includes security validation for file content.
    
    Args:
        filepath: Path to the file to extract text from
        
    Returns:
        Tuple of (extracted_text, file_extension)
    """
    if not os.path.exists(filepath):
        logger.error(f"File not found: {filepath}")
        return None, ""
    
    # Get file extension
    _, ext = os.path.splitext(filepath.lower())
    
    # Try to detect MIME type for additional validation
    mime_type = _get_file_mime_type(filepath)
    if mime_type and not _is_supported_mime_type(mime_type):
        logger.warning(f"Unsupported MIME type {mime_type} for file {filepath}")
        return None, ext
    
    # Extract text based on file extension
    if ext == '.txt':
        content = get_text_from_txt(filepath)
    elif ext == '.pdf':
        content = get_text_from_pdf(filepath)
    elif ext == '.docx':
        content = get_text_from_docx(filepath)
    elif ext == '.xlsx':
        content = get_text_from_xlsx(filepath)
    elif ext == '.pptx':
        content = get_text_from_pptx(filepath)
    elif ext == '.csv':
        content = get_text_from_csv(filepath)
    elif ext == '.doc':
        content = get_text_from_doc(filepath)
    else:
        logger.warning(f"Unsupported file extension: {ext}")
        return None, ext
    
    # Final content validation
    if content and not _validate_file_content(filepath, content):
        logger.warning(f"Final content validation failed for {filepath}")
        return None, ext
    
    return content, ext

# Test function for standalone testing
if __name__ == '__main__':
    print("--- Running file_parsers.py standalone tests ---")
    
    # Test with a simple text file
    test_content = "This is a test document with some content."
    test_file = "test_document.txt"
    
    try:
        with open(test_file, 'w', encoding='utf-8') as f:
            f.write(test_content)
        
        text, ext = get_text_from_file(test_file)
        print(f"Test file ({ext}): {text}")
        
    finally:
        if os.path.exists(test_file):
            os.remove(test_file)
    
    print("--- Test complete ---")